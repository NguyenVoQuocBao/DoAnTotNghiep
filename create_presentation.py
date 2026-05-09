from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# Colors from DESIGN.md
C_BG = RGBColor(0x09, 0x09, 0x0b)        # #09090b surface
C_SURFACE = RGBColor(0x11, 0x11, 0x13)   # #111113
C_PRIMARY = RGBColor(0x0C, 0x5C, 0xAB)   # #0C5CAB
C_TEXT = RGBColor(0xfa, 0xfa, 0xfa)      # #fafafa
C_MUTED = RGBColor(0xa1, 0xa1, 0xaa)     # #a1a1aa
C_SUCCESS = RGBColor(0x10, 0xb9, 0x81)   # #10b981
C_WARNING = RGBColor(0xf5, 0x9e, 0x0b)   # #f59e0b
C_DANGER = RGBColor(0xef, 0x44, 0x44)    # #ef4444
C_WHITE = RGBColor(0xff, 0xff, 0xff)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def add_bg(slide, color=C_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=18, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = 'IBM Plex Sans'
    return txBox

def add_bullet_slide(prs, title, bullets, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    # Left accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.06), H, C_PRIMARY)
    # Title
    add_text(slide, title, Inches(0.4), Inches(0.35), Inches(12.5), Inches(0.8),
             size=28, bold=True, color=C_TEXT)
    # Divider
    add_rect(slide, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.03), C_PRIMARY)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.5),
                 size=14, color=C_MUTED)
    y = Inches(1.75) if not subtitle else Inches(1.85)
    for b in bullets:
        if b.startswith('##'):
            add_text(slide, b[2:].strip(), Inches(0.4), y, Inches(12.5), Inches(0.5),
                     size=16, bold=True, color=C_PRIMARY)
            y += Inches(0.5)
        else:
            add_rect(slide, Inches(0.4), y + Inches(0.12), Inches(0.06), Inches(0.06), C_PRIMARY)
            add_text(slide, b, Inches(0.6), y, Inches(12.2), Inches(0.45),
                     size=15, color=C_TEXT)
            y += Inches(0.48)
    return slide

# ── SLIDE 1: COVER ──────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), C_PRIMARY)
add_rect(slide, Inches(0), H - Inches(0.06), W, Inches(0.06), C_PRIMARY)
# Big title
add_text(slide, 'XÂY DỰNG MÔ HÌNH DỰ ĐOÁN', Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.0),
         size=36, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
add_text(slide, 'MỨC ĐỘ HOÀN THÀNH CỦA SINH VIÊN', Inches(0.8), Inches(2.4), Inches(11.5), Inches(1.0),
         size=36, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(4.5), Inches(3.5), Inches(4.3), Inches(0.04), C_PRIMARY)
add_text(slide, 'Đồ án tốt nghiệp', Inches(0.8), Inches(3.7), Inches(11.5), Inches(0.5),
         size=16, color=C_MUTED, align=PP_ALIGN.CENTER)
add_text(slide, 'Sinh viên: Nguyễn Minh Phúc  |  MSSV: 122001605', Inches(0.8), Inches(4.3), Inches(11.5), Inches(0.5),
         size=15, color=C_TEXT, align=PP_ALIGN.CENTER)
add_text(slide, 'Giảng viên hướng dẫn: Nguyễn Võ Quốc Bảo', Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.5),
         size=15, color=C_TEXT, align=PP_ALIGN.CENTER)
add_text(slide, '2025', Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
         size=14, color=C_MUTED, align=PP_ALIGN.CENTER)

# ── SLIDE 2: MỤC LỤC ────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), Inches(0.06), H, C_PRIMARY)
add_text(slide, 'Nội dung trình bày', Inches(0.4), Inches(0.35), Inches(12.5), Inches(0.8),
         size=28, bold=True, color=C_TEXT)
add_rect(slide, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.03), C_PRIMARY)
items = [
    ('01', 'Giới thiệu đề tài & Mục tiêu'),
    ('02', 'Tổng quan hệ thống'),
    ('03', 'Dữ liệu & Tiền xử lý'),
    ('04', 'Mô hình Machine Learning'),
    ('05', 'Kết quả & Đánh giá'),
    ('06', 'Demo hệ thống'),
    ('07', 'Kết luận & Hướng phát triển'),
]
for i, (num, text) in enumerate(items):
    y = Inches(1.3) + Inches(i * 0.72)
    add_rect(slide, Inches(0.4), y, Inches(0.55), Inches(0.45), C_PRIMARY)
    add_text(slide, num, Inches(0.4), y, Inches(0.55), Inches(0.45),
             size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, text, Inches(1.1), y + Inches(0.05), Inches(11), Inches(0.4),
             size=16, color=C_TEXT)

# ── SLIDE 3: GIỚI THIỆU ─────────────────────────────────────
add_bullet_slide(prs, '01. Giới thiệu đề tài', [
    '## Vấn đề thực tiễn',
    'Tỷ lệ sinh viên không hoàn thành môn học ngày càng tăng',
    'Giảng viên khó theo dõi từng sinh viên trong lớp đông',
    'Thiếu công cụ phân tích và cảnh báo sớm',
    '## Mục tiêu đề tài',
    'Xây dựng mô hình dự đoán mức độ hoàn thành học tập',
    'Phân tích hành vi học tập và đưa ra gợi ý cá nhân hóa',
    'Xây dựng hệ thống web hỗ trợ giảng viên quản lý lớp học',
])

# ── SLIDE 4: TỔNG QUAN HỆ THỐNG ─────────────────────────────
add_bullet_slide(prs, '02. Tổng quan hệ thống', [
    '## Kiến trúc',
    'Backend: Python Flask + REST API',
    'ML Engine: scikit-learn (Random Forest + Neural Network)',
    'Storage: JSON file-based (không cần database)',
    'Frontend: HTML/CSS/JS + Chart.js',
    '## Người dùng',
    'Sinh viên: xem phân tích cá nhân, lịch sử học tập',
    'Giảng viên: quản lý lớp, xuất báo cáo, import dữ liệu',
])

# ── SLIDE 5: DỮ LIỆU ────────────────────────────────────────
add_bullet_slide(prs, '03. Dữ liệu & Tiền xử lý', [
    '## Dữ liệu đầu vào (5 đặc trưng)',
    'Điểm kiểm tra (test_scores)',
    'Tỷ lệ hoàn thành bài tập (assignment_completion_rate)',
    'Tần suất truy cập hệ thống (access_count)',
    'Thời gian học tập (study_time_minutes)',
    'Thời điểm nộp bài (submission_timing_score)',
    '## Chuẩn hóa (DataNormalizer)',
    'Tất cả đặc trưng được chuẩn hóa về thang 1–10',
    'Công thức: normalized = (value / max_value) × 10, clamp [1, 10]',
])

# ── SLIDE 6: MÔ HÌNH ML ─────────────────────────────────────
add_bullet_slide(prs, '04. Mô hình Machine Learning', [
    '## Random Forest Classifier',
    '100 cây quyết định, 5-Fold Stratified Cross Validation',
    'Độ chính xác CV trung bình: ~97%',
    '## Neural Network (MLPClassifier)',
    'Validation accuracy: ~90%',
    '## Ensemble',
    'Kết hợp RF + NN bằng averaging xác suất',
    '## Công thức tính điểm hoàn thành',
    'Điểm = KT×30% + BT×25% + Truy cập×15% + TG học×15% + Nộp bài×15%',
])

# ── SLIDE 7: KẾT QUẢ ────────────────────────────────────────
add_bullet_slide(prs, '05. Kết quả & Đánh giá', [
    '## Phân loại mức độ hoàn thành',
    'Cao (≥ 7.0 điểm)  |  Trung bình (4.0–7.0)  |  Thấp (< 4.0)',
    '## Phân nhóm hành vi',
    'Nhóm 1: Thụ động / ít tham gia',
    'Nhóm 2: Học tập chưa ổn định / có xu hướng trì hoãn',
    'Nhóm 3: Học tập tích cực và có chiến lược',
    '## Hiệu suất mô hình',
    'Random Forest CV accuracy: 97.4% ± 1.6%',
    'Precision / Recall / F1 = 1.000 trên tập huấn luyện',
])

# ── SLIDE 8: DEMO ────────────────────────────────────────────
def add_image_slide(prs, title, img_path, caption=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.06), H, C_PRIMARY)
    add_text(slide, title, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.6),
             size=22, bold=True, color=C_TEXT)
    add_rect(slide, Inches(0.4), Inches(0.85), Inches(12.5), Inches(0.03), C_PRIMARY)
    if os.path.exists(img_path):
        img_top = Inches(0.95)
        img_h = Inches(5.8) if not caption else Inches(5.5)
        slide.shapes.add_picture(img_path, Inches(0.4), img_top, width=Inches(12.5), height=img_h)
    if caption:
        add_text(slide, caption, Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.4),
                 size=12, color=C_MUTED, align=PP_ALIGN.CENTER, italic=True)
    return slide

add_image_slide(prs, '06. Demo — Dashboard Giảng viên',
    'screenshot_teacher_dashboard.png',
    'Tổng quan lớp học: 196 sinh viên, phân bố mức độ hoàn thành và hành vi học tập')

add_image_slide(prs, '06. Demo — Thống kê & Biểu đồ',
    'screenshot_statistics.png',
    'Biểu đồ phân bố mức độ hoàn thành, nhóm hành vi và histogram điểm số')

add_image_slide(prs, '06. Demo — Báo cáo Tổng kết Lớp',
    'screenshot_class_report.png',
    'Báo cáo đầy đủ: Top 10 xuất sắc, danh sách sinh viên cần hỗ trợ')

add_image_slide(prs, '06. Demo — Dashboard Sinh viên',
    'screenshot_student_dashboard.png',
    'Sinh viên xem mức độ hoàn thành, radar chart kỹ năng và dự đoán AI')

add_image_slide(prs, '06. Demo — Phân tích Chi tiết Sinh viên',
    'screenshot_student_analysis.png',
    'Biểu đồ năng lực, dự đoán RF + NN, xác suất và gợi ý học tập cá nhân hóa')

add_image_slide(prs, '06. Demo — Chi tiết Sinh viên (Giảng viên)',
    'screenshot_student_detail.png',
    'Giảng viên xem chi tiết từng sinh viên: kết quả phân tích, gợi ý, dữ liệu chuẩn hóa')

add_image_slide(prs, '06. Demo — So sánh Sinh viên',
    'screenshot_compare.png',
    'So sánh 2-4 sinh viên cạnh nhau: radar chart, bar chart và bảng tổng hợp')

add_image_slide(prs, '05. Kết quả — Giải thích Mô hình',
    'screenshot_model.png',
    'Feature importance, Confusion Matrix, CV scores 97.5%, Precision/Recall/F1')

# ── SLIDE 9: KẾT LUẬN ───────────────────────────────────────
add_bullet_slide(prs, '07. Kết luận & Hướng phát triển', [
    '## Kết quả đạt được',
    'Xây dựng thành công mô hình RF với độ chính xác cao',
    'Hệ thống web đầy đủ tính năng cho giảng viên và sinh viên',
    'Tích hợp Advanced Analytics: Clustering, Anomaly Detection',
    '## Hạn chế',
    'Dữ liệu lưu trữ JSON — chưa phù hợp cho quy mô lớn',
    'Chưa có tính năng thông báo real-time',
    '## Hướng phát triển',
    'Chuyển sang database (PostgreSQL/MongoDB)',
    'Tích hợp LMS (Moodle, Canvas) để lấy dữ liệu tự động',
])

# ── SLIDE 10: CẢM ƠN ────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), C_PRIMARY)
add_rect(slide, Inches(0), H - Inches(0.06), W, Inches(0.06), C_PRIMARY)
add_text(slide, 'Xin cảm ơn!', Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.2),
         size=48, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
add_text(slide, 'Kính mong nhận được ý kiến đóng góp từ Hội đồng', Inches(0.8), Inches(3.5), Inches(11.5), Inches(0.6),
         size=18, color=C_MUTED, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(4.5), Inches(4.2), Inches(4.3), Inches(0.04), C_PRIMARY)
add_text(slide, 'Nguyễn Minh Phúc — 122001605', Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.5),
         size=15, color=C_TEXT, align=PP_ALIGN.CENTER)

output = 'Bao_cao_tot_nghiep_NguyenMinhPhuc.pptx'
prs.save(output)
print(f'Done: {output}')
